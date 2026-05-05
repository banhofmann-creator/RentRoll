from __future__ import annotations

import shutil
import time
from dataclasses import asdict
from datetime import date
from io import BytesIO
from pathlib import Path
from uuid import uuid4

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.config import settings
from app.core.kpi_catalog import KPI_CATALOG, format_value
from app.core.pptx_patcher import Mapping, apply_token_mappings
from app.models.database import CsvUpload, PptxRefreshJob, RawRentRoll, ReportingPeriod
from app.parsers.pptx_ingestor import find_token_candidates, ingest_pptx


@pytest.fixture
def pptx_upload_dir(monkeypatch):
    base = Path.cwd() / ".tmp_pptx_tests"
    base.mkdir(exist_ok=True)
    root = base / f"tmp_{uuid4().hex}"
    root.mkdir()
    old_upload_dir = settings.upload_dir
    monkeypatch.setattr(settings, "upload_dir", str(root))
    try:
        yield root
    finally:
        monkeypatch.setattr(settings, "upload_dir", old_upload_dir)
        shutil.rmtree(root, ignore_errors=True)


def _make_synthetic_deck(
    tokens: list[str],
    table_tokens: list[str] | None = None,
    surrounding: bool = False,
) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    top = Inches(0.5)
    for token in tokens:
        box = slide.shapes.add_textbox(Inches(0.5), top, Inches(5), Inches(0.5))
        paragraph = box.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = f"Portfolio rent: {token} EUR" if surrounding else token
        top += Inches(0.6)

    table_tokens = table_tokens or []
    if table_tokens:
        table_shape = slide.shapes.add_table(2, 2, Inches(0.5), Inches(3), Inches(5), Inches(1.2))
        table = table_shape.table
        for idx, token in enumerate(table_tokens[:4]):
            row = idx // 2
            col = idx % 2
            paragraph = table.cell(row, col).text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = token

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _deck_text(file_bytes: bytes) -> str:
    prs = Presentation(BytesIO(file_bytes))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text_frame.text)
    return " ".join(parts)


def _create_period_with_data(db, status="draft") -> ReportingPeriod:
    upload = CsvUpload(
        filename="test.csv",
        status="complete",
        stichtag=date(2025, 3, 31),
        row_count=2,
        data_row_count=2,
        summary_row_count=0,
    )
    db.add(upload)
    db.commit()
    db.refresh(upload)

    db.add_all(
        [
            RawRentRoll(
                upload_id=upload.id,
                row_number=1,
                row_type="data",
                fund="GLIF",
                property_id="1001",
                tenant_name="Tenant A",
                unit_type="Halle",
                area_sqm=1000,
                annual_net_rent=120000,
                monthly_net_rent=10000,
            ),
            RawRentRoll(
                upload_id=upload.id,
                row_number=2,
                row_type="data",
                fund="GLIF",
                property_id="1001",
                tenant_name="LEERSTAND",
                unit_type="Halle",
                area_sqm=100,
                annual_net_rent=0,
                monthly_net_rent=0,
            ),
        ]
    )
    period = ReportingPeriod(stichtag=date(2025, 3, 31), upload_id=upload.id, status=status)
    db.add(period)
    db.commit()
    db.refresh(period)
    return period


def _wait_for_status(client, job_id: int, status: str = "proposed") -> dict:
    for _ in range(20):
        response = client.get(f"/api/pptx/{job_id}")
        assert response.status_code == 200
        data = response.json()
        if data["status"] == status or data["status"] == "error":
            return data
        time.sleep(0.1)
    raise AssertionError(f"PPTX refresh job {job_id} did not reach {status}")


def test_kpi_catalog_has_all_snapshot_keys():
    expected = {
        "total_rent",
        "total_area",
        "vacant_area",
        "vacancy_rate",
        "tenant_count",
        "property_count",
        "fair_value",
        "total_debt",
        "wault_avg",
    }
    assert set(KPI_CATALOG) == expected


def test_format_value_money_eur_millions():
    assert format_value(12_500_000, "money_eur_millions") == "12,5 M€"


def test_format_value_percent():
    assert format_value(5.32, "percent") == "5,32 %"


def test_format_value_percent_below_one():
    assert format_value(0.5, "percent") == "0,5 %"


def test_format_value_integer():
    assert format_value(17, "integer") == "17"


def test_ingest_finds_text_frame_tokens():
    elements = ingest_pptx(_make_synthetic_deck(["{{total_rent}}"]))
    tokens, unknown_tokens = find_token_candidates(elements)

    text_tokens = [token for token in tokens if token.address.kind == "text_frame"]
    assert unknown_tokens == []
    assert len(text_tokens) == 1
    assert text_tokens[0].kpi_id == "total_rent"


def test_ingest_finds_table_cell_tokens():
    elements = ingest_pptx(_make_synthetic_deck([], table_tokens=["{{total_area}}"]))
    tokens, _ = find_token_candidates(elements)

    table_tokens = [token for token in tokens if token.address.kind == "table_cell"]
    assert len(table_tokens) == 1
    assert table_tokens[0].kpi_id == "total_area"
    assert table_tokens[0].address.row is not None
    assert table_tokens[0].address.col is not None


def test_ingest_unknown_token():
    elements = ingest_pptx(_make_synthetic_deck(["{{not_a_kpi}}"]))
    tokens, unknown_tokens = find_token_candidates(elements)

    assert tokens == []
    assert unknown_tokens == ["not_a_kpi"]


def test_patcher_replaces_token_in_text_frame():
    source = _make_synthetic_deck(["{{total_rent}}"])
    tokens, _ = find_token_candidates(ingest_pptx(source))

    patched, changes = apply_token_mappings(
        source,
        [Mapping(tokens[0].address, "{{total_rent}}", "12,5 M€")],
    )

    assert changes[0].success is True
    assert "12,5 M€" in _deck_text(patched)
    assert "{{total_rent}}" not in _deck_text(patched)


def test_patcher_replaces_token_in_table_cell():
    source = _make_synthetic_deck([], table_tokens=["{{total_area}}"])
    tokens, _ = find_token_candidates(ingest_pptx(source))

    patched, changes = apply_token_mappings(
        source,
        [Mapping(tokens[0].address, "{{total_area}}", "1.100 m²")],
    )

    assert changes[0].success is True
    assert "1.100 m²" in _deck_text(patched)


def test_patcher_preserves_other_text():
    source = _make_synthetic_deck(["{{total_rent}}"], surrounding=True)
    tokens, _ = find_token_candidates(ingest_pptx(source))

    patched, _ = apply_token_mappings(
        source,
        [Mapping(tokens[0].address, "{{total_rent}}", "12,5 M€")],
    )

    text = _deck_text(patched)
    assert "Portfolio rent:" in text
    assert "EUR" in text
    assert "12,5 M€" in text


def test_api_upload_and_status_flow(client, pptx_upload_dir):
    deck = _make_synthetic_deck(["{{total_rent}}"])

    response = client.post(
        "/api/pptx/upload",
        files={"file": ("deck.pptx", deck, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )

    assert response.status_code == 200
    data = response.json()
    assert data["status"] == "uploaded"

    job = _wait_for_status(client, data["id"])
    assert job["status"] == "proposed"
    assert job["proposals"]["tokens"] == ["total_rent"]
    assert job["proposals_json"]["unknown_tokens"] == []


def test_api_apply_complete(client, db, pptx_upload_dir):
    period = _create_period_with_data(db)
    deck = _make_synthetic_deck(["{{total_rent}}"])
    upload_response = client.post(
        "/api/pptx/upload",
        files={"file": ("deck.pptx", deck, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    job_id = upload_response.json()["id"]
    _wait_for_status(client, job_id)

    response = client.post(f"/api/pptx/{job_id}/apply", json={"period_id": period.id, "mappings": None})

    assert response.status_code == 200
    data = response.json()
    assert data["status"] == "complete"
    assert data["output_filename"] == "deck-refreshed.pptx"
    assert (pptx_upload_dir / data["output_blob_path"]).exists()


def test_api_apply_records_period_status_for_draft(client, db, pptx_upload_dir):
    period = _create_period_with_data(db, status="draft")
    deck = _make_synthetic_deck(["{{total_rent}}"])
    upload_response = client.post(
        "/api/pptx/upload",
        files={"file": ("deck.pptx", deck, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    job_id = upload_response.json()["id"]
    _wait_for_status(client, job_id)

    response = client.post(f"/api/pptx/{job_id}/apply", json={"period_id": period.id, "mappings": None})

    assert response.status_code == 200
    assert response.json()["period_status_at_refresh"] == "draft"


def test_api_apply_unknown_kpi_returns_error(client, db, pptx_upload_dir):
    period = _create_period_with_data(db)
    deck = _make_synthetic_deck(["{{not_a_kpi}}"])
    elements = ingest_pptx(deck)
    address = asdict(elements[0].address)

    job_dir = pptx_upload_dir / "pptx_refresh" / "99"
    job_dir.mkdir(parents=True)
    (job_dir / "source.pptx").write_bytes(deck)
    job = PptxRefreshJob(
        id=99,
        original_filename="deck.pptx",
        original_blob_path=str(Path("pptx_refresh") / "99" / "source.pptx"),
        status="proposed",
        proposals_json={
            "mode": "token",
            "tokens": [
                {
                    "address": address,
                    "kpi_id": "not_a_kpi",
                    "full_text": "{{not_a_kpi}}",
                    "span": [0, 13],
                }
            ],
            "unknown_tokens": [],
        },
    )
    db.add(job)
    db.commit()

    response = client.post(f"/api/pptx/{job.id}/apply", json={"period_id": period.id, "mappings": None})

    assert response.status_code == 200
    data = response.json()
    assert data["status"] == "error"
    assert "Unknown KPI token" in data["error_message"]


def test_api_download_returns_pptx(client, db, pptx_upload_dir):
    period = _create_period_with_data(db)
    deck = _make_synthetic_deck(["{{total_rent}}"])
    upload_response = client.post(
        "/api/pptx/upload",
        files={"file": ("deck.pptx", deck, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    job_id = upload_response.json()["id"]
    _wait_for_status(client, job_id)
    apply_response = client.post(f"/api/pptx/{job_id}/apply", json={"period_id": period.id, "mappings": None})
    assert apply_response.json()["status"] == "complete"

    response = client.get(f"/api/pptx/{job_id}/download")

    assert response.status_code == 200
    assert response.headers["content-type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert 'filename="deck-refreshed.pptx"' in response.headers["content-disposition"]
    assert Presentation(BytesIO(response.content))
