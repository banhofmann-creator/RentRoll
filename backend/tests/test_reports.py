"""Tests for reports API — PPTX slide generation."""
from datetime import date
from io import BytesIO

import pytest
from pptx import Presentation

from app.models.database import CsvUpload, PropertyMaster, RawRentRoll


def _setup_upload(db, stichtag=date(2025, 3, 31)):
    upload = CsvUpload(filename="test.csv", status="complete", stichtag=stichtag)
    db.add(upload)
    db.commit()
    db.refresh(upload)

    db.add(RawRentRoll(
        upload_id=upload.id, row_number=1, row_type="data",
        fund="GLIF", property_id="7042",
        tenant_name="Tenant A", unit_type="Halle",
        area_sqm=5000, annual_net_rent=120000,
        lease_end_agreed=date(2027, 6, 30),
    ))
    db.add(RawRentRoll(
        upload_id=upload.id, row_number=2, row_type="data",
        fund="GLIF", property_id="7042",
        tenant_name="Tenant B", unit_type="Büro",
        area_sqm=500, annual_net_rent=30000,
        lease_end_agreed=date(2029, 12, 31),
    ))
    db.add(RawRentRoll(
        upload_id=upload.id, row_number=3, row_type="data",
        fund="GLIF", property_id="7042",
        tenant_name="LEERSTAND", unit_type="Halle",
        area_sqm=1000, annual_net_rent=0,
    ))
    db.add(RawRentRoll(
        upload_id=upload.id, row_number=4, row_type="data",
        fund="GLIFPLUSII", property_id="1001",
        tenant_name="Tenant C", unit_type="Halle",
        area_sqm=8000, annual_net_rent=200000,
        lease_end_agreed=date(2030, 3, 31),
    ))
    db.add(RawRentRoll(
        upload_id=upload.id, row_number=5, row_type="property_summary",
        fund="7042 - Almere",
        area_sqm=6500, annual_net_rent=150000, wault=4.5,
        property_id="7042",
    ))
    db.commit()

    db.add(PropertyMaster(property_id="7042", city="Almere", country="NL"))
    db.add(PropertyMaster(property_id="1001", city="Hamburg", country="DE"))
    db.commit()

    return upload


def _read_pptx(response) -> Presentation:
    assert response.status_code == 200
    content_type = response.headers.get("content-type", "")
    assert "presentation" in content_type
    return Presentation(BytesIO(response.content))


def test_property_factsheet(client, db):
    upload = _setup_upload(db)
    resp = client.get(f"/api/reports/property-factsheet?upload_id={upload.id}&property_id=7042")
    prs = _read_pptx(resp)
    assert len(prs.slides) >= 2


def test_property_factsheet_not_found(client, db):
    upload = _setup_upload(db)
    resp = client.get(f"/api/reports/property-factsheet?upload_id={upload.id}&property_id=9999")
    assert resp.status_code == 404


def test_portfolio_overview(client, db):
    upload = _setup_upload(db)
    resp = client.get(f"/api/reports/portfolio-overview?upload_id={upload.id}")
    prs = _read_pptx(resp)
    assert len(prs.slides) >= 3


def test_lease_expiry_profile(client, db):
    upload = _setup_upload(db)
    resp = client.get(f"/api/reports/lease-expiry?upload_id={upload.id}")
    prs = _read_pptx(resp)
    assert len(prs.slides) >= 2


def test_fund_summary(client, db):
    upload = _setup_upload(db)
    resp = client.get(f"/api/reports/fund-summary?upload_id={upload.id}&fund=GLIF")
    prs = _read_pptx(resp)
    assert len(prs.slides) >= 2


def test_fund_summary_not_found(client, db):
    upload = _setup_upload(db)
    resp = client.get(f"/api/reports/fund-summary?upload_id={upload.id}&fund=NONEXISTENT")
    assert resp.status_code == 404


def test_available_funds(client, db):
    upload = _setup_upload(db)
    resp = client.get(f"/api/reports/available-funds?upload_id={upload.id}")
    assert resp.status_code == 200
    funds = resp.json()
    assert "GLIF" in funds
    assert "GLIFPLUSII" in funds


def test_available_properties(client, db):
    upload = _setup_upload(db)
    resp = client.get(f"/api/reports/available-properties?upload_id={upload.id}")
    assert resp.status_code == 200
    props = resp.json()
    assert "7042" in props
    assert "1001" in props


def test_upload_not_found(client):
    resp = client.get("/api/reports/portfolio-overview?upload_id=9999")
    assert resp.status_code == 404


def test_factsheet_has_correct_content(client, db):
    upload = _setup_upload(db)
    resp = client.get(f"/api/reports/property-factsheet?upload_id={upload.id}&property_id=7042")
    prs = _read_pptx(resp)

    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text + " "

    assert "7042" in all_text
    assert "Almere" in all_text or "NL" in all_text
