from datetime import date
from io import BytesIO
from pathlib import Path
import shutil
from uuid import uuid4
from zipfile import ZipFile

import openpyxl
import pytest
from pptx import Presentation

import app.channels.registry as channel_registry
from app.channels.local_filesystem import LocalFilesystemChannel
from app.channels.registry import get_channel, list_channels
from app.channels.base import ExportFile, ExportMetadata
from app.core.investor_pack import generate_investor_pack
from app.models.database import CsvUpload, PropertyMaster, RawRentRoll, ReportingPeriod


@pytest.fixture
def tmp_path():
    root = Path.cwd() / ".tmp_export_tests"
    root.mkdir(exist_ok=True)
    path = root / f"tmp_{uuid4().hex}"
    path.mkdir()
    try:
        yield path
    finally:
        shutil.rmtree(path, ignore_errors=True)


def _create_reporting_period(db) -> ReportingPeriod:
    upload = CsvUpload(
        filename="test.csv",
        status="complete",
        stichtag=date(2025, 3, 31),
        fund_label="TESTFUND",
    )
    db.add(upload)
    db.commit()
    db.refresh(upload)

    db.add_all(
        [
            RawRentRoll(
                upload_id=upload.id,
                row_number=1,
                row_type="data",
                fund="TESTFUND",
                property_id="1001",
                tenant_name="TestTenant",
                unit_type="Halle",
                area_sqm=100.0,
                annual_net_rent=12000.0,
                monthly_net_rent=1000.0,
                lease_end_agreed=date(2027, 6, 30),
            ),
            RawRentRoll(
                upload_id=upload.id,
                row_number=2,
                row_type="data",
                fund="TESTFUND",
                property_id="1001",
                tenant_name="LEERSTAND",
                unit_type="Halle",
                area_sqm=25.0,
                annual_net_rent=0.0,
                monthly_net_rent=0.0,
            ),
            RawRentRoll(
                upload_id=upload.id,
                row_number=3,
                row_type="data",
                fund="SECONDFUND",
                property_id="2002",
                tenant_name="SecondTenant",
                unit_type="Buro",
                area_sqm=200.0,
                annual_net_rent=24000.0,
                monthly_net_rent=2000.0,
                lease_end_agreed=date(2028, 12, 31),
            ),
        ]
    )
    db.add_all(
        [
            PropertyMaster(property_id="1001", city="Hamburg", country="DE"),
            PropertyMaster(property_id="2002", city="Berlin", country="DE"),
        ]
    )
    db.commit()

    period = ReportingPeriod(
        stichtag=date(2025, 3, 31),
        upload_id=upload.id,
        status="draft",
    )
    db.add(period)
    db.commit()
    db.refresh(period)
    return period


def test_list_channels_includes_local_filesystem():
    channels = {channel["name"]: channel for channel in list_channels()}
    assert channels["local_filesystem"] == {
        "name": "local_filesystem",
        "description": LocalFilesystemChannel.description,
    }


def test_get_channel_returns_local_filesystem_instance():
    channel = get_channel("local_filesystem")
    assert isinstance(channel, LocalFilesystemChannel)


def test_get_channel_unknown_raises_value_error():
    with pytest.raises(ValueError, match="Unknown channel: nonexistent"):
        get_channel("nonexistent")


def test_local_filesystem_channel_connection_and_push(tmp_path):
    channel = LocalFilesystemChannel(tmp_path / "exports")

    assert channel.test_connection() is True

    files = [
        ExportFile(
            filename="investor-pack.zip",
            content=b"zip-bytes",
            file_type="zip",
            category="investor_pack",
        ),
        ExportFile(
            filename="portfolio_overview.pptx",
            content=b"pptx-bytes",
            file_type="pptx",
            category="portfolio_overview",
        ),
    ]
    metadata = ExportMetadata(
        stichtag=date(2025, 3, 31),
        fund="TESTFUND",
        properties=["1001"],
        reporting_period_id=1,
    )

    result = channel.push(files, metadata)
    destination = tmp_path / "exports" / "TESTFUND" / "2025-03-31"

    assert result.success is True
    assert result.files_pushed == 2
    assert result.destination == str(destination)
    assert (destination / "investor-pack.zip").read_bytes() == b"zip-bytes"
    assert (destination / "portfolio_overview.pptx").read_bytes() == b"pptx-bytes"


def test_generate_investor_pack_returns_valid_zip_bytes(db):
    period = _create_reporting_period(db)

    zip_bytes, filename, export_files = generate_investor_pack(db, period.id)

    assert filename == "investor_pack_all_funds_2025-03-31.zip"
    assert len(export_files) == 7

    with ZipFile(BytesIO(zip_bytes)) as archive:
        names = set(archive.namelist())
        assert "BVI_all_funds_2025-03-31_DRAFT.xlsx" in names
        assert "portfolio_overview_2025-03-31.pptx" in names
        assert "lease_expiry_profile_2025-03-31.pptx" in names

        workbook = openpyxl.load_workbook(
            BytesIO(archive.read("BVI_all_funds_2025-03-31_DRAFT.xlsx"))
        )
        assert "Z1_Tenants_Leases" in workbook.sheetnames
        assert "G2_Property_data" in workbook.sheetnames

        portfolio_overview = Presentation(
            BytesIO(archive.read("portfolio_overview_2025-03-31.pptx"))
        )
        lease_expiry = Presentation(
            BytesIO(archive.read("lease_expiry_profile_2025-03-31.pptx"))
        )
        assert len(portfolio_overview.slides) >= 2
        assert len(lease_expiry.slides) >= 2


def test_generate_investor_pack_fund_filter_only_includes_requested_fund(db):
    period = _create_reporting_period(db)

    _, filename, export_files = generate_investor_pack(db, period.id, fund="TESTFUND")

    file_names = {export_file.filename for export_file in export_files}
    assert filename == "investor_pack_TESTFUND_2025-03-31.zip"
    assert len(file_names) == 5
    assert "fund_summary_TESTFUND_2025-03-31.pptx" in file_names
    assert "factsheet_1001_2025-03-31.pptx" in file_names
    assert "fund_summary_SECONDFUND_2025-03-31.pptx" not in file_names
    assert "factsheet_2002_2025-03-31.pptx" not in file_names


def test_generate_investor_pack_period_not_found_raises_value_error(db):
    with pytest.raises(ValueError, match="Reporting period not found"):
        generate_investor_pack(db, 999999)


def test_export_channels_endpoint_returns_local_filesystem(client):
    response = client.get("/api/export/channels")

    assert response.status_code == 200
    assert {
        "name": "local_filesystem",
        "description": LocalFilesystemChannel.description,
    } in response.json()


def test_export_investor_pack_endpoint_returns_zip(client, db):
    period = _create_reporting_period(db)

    response = client.post(f"/api/export/investor-pack?period_id={period.id}")

    assert response.status_code == 200
    assert response.headers["content-type"] == "application/zip"
    assert 'filename="investor_pack_all_funds_2025-03-31.zip"' in response.headers["content-disposition"]

    with ZipFile(BytesIO(response.content)) as archive:
        names = set(archive.namelist())
        assert "BVI_all_funds_2025-03-31_DRAFT.xlsx" in names
        assert "portfolio_overview_2025-03-31.pptx" in names
        assert "lease_expiry_profile_2025-03-31.pptx" in names


def test_export_investor_pack_preview_endpoint_returns_manifest(client, db):
    period = _create_reporting_period(db)

    response = client.post(
        f"/api/export/investor-pack/preview?period_id={period.id}&fund=TESTFUND"
    )

    assert response.status_code == 200
    data = response.json()
    assert data["filename"] == "investor_pack_TESTFUND_2025-03-31.zip"
    assert data["file_count"] == 5
    assert "BVI_TESTFUND_2025-03-31_DRAFT.xlsx" in data["files"]
    assert "fund_summary_TESTFUND_2025-03-31.pptx" in data["files"]
    assert "factsheet_1001_2025-03-31.pptx" in data["files"]


def test_export_investor_pack_endpoint_returns_404_for_missing_period(client):
    response = client.post("/api/export/investor-pack?period_id=999999")

    assert response.status_code == 404


def test_export_push_endpoint_returns_success_for_local_filesystem(client, db, tmp_path, monkeypatch):
    period = _create_reporting_period(db)

    class TempLocalFilesystemChannel(LocalFilesystemChannel):
        def __init__(self):
            super().__init__(tmp_path / "exports")

    monkeypatch.setitem(channel_registry._CHANNELS, "local_filesystem", TempLocalFilesystemChannel)

    response = client.post(
        "/api/export/push",
        json={
            "period_id": period.id,
            "channel": "local_filesystem",
            "fund": "TESTFUND",
        },
    )

    assert response.status_code == 200
    data = response.json()
    destination = tmp_path / "exports" / "TESTFUND" / "2025-03-31"
    assert data["success"] is True
    assert data["files_pushed"] == 5
    assert Path(data["destination"]) == destination
    assert (destination / "BVI_TESTFUND_2025-03-31_DRAFT.xlsx").exists()
    assert (destination / "fund_summary_TESTFUND_2025-03-31.pptx").exists()


def test_export_push_endpoint_returns_400_for_invalid_channel(client, db):
    period = _create_reporting_period(db)

    response = client.post(
        "/api/export/push",
        json={"period_id": period.id, "channel": "invalid_channel"},
    )

    assert response.status_code == 400
    assert response.json()["detail"] == "Unknown channel: invalid_channel"
