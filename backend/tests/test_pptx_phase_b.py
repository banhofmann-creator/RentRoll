"""Tests for RR-1 Phase B (AI-driven PPTX KPI resolver).

These tests do not call the live Anthropic API; they inject a fake client
through ``app.api.pptx_refresh.set_ai_client_override``.
"""
from __future__ import annotations

import json
import shutil
import time
from dataclasses import dataclass
from datetime import date
from io import BytesIO
from pathlib import Path
from uuid import uuid4

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.api.pptx_refresh import set_ai_client_override
from app.config import settings
from app.core.pptx_kpi_resolver import (
    Candidate,
    collect_candidates,
    resolve_with_ai,
)
from app.models.database import CsvUpload, RawRentRoll, ReportingPeriod
from app.parsers.pptx_ingestor import ingest_pptx


# ── Fixtures ──────────────────────────────────────────────────────────


@pytest.fixture
def pptx_upload_dir(monkeypatch):
    base = Path.cwd() / ".tmp_pptx_b_tests"
    base.mkdir(exist_ok=True)
    root = base / f"tmp_{uuid4().hex}"
    root.mkdir()
    old_upload_dir = settings.upload_dir
    monkeypatch.setattr(settings, "upload_dir", str(root))
    try:
        yield root
    finally:
        monkeypatch.setattr(settings, "upload_dir", old_upload_dir)
        shutil.rmtree(root, ignore_errors=True)


# ── Fake Anthropic client ─────────────────────────────────────────────


@dataclass
class _FakeBlock:
    type: str
    text: str


@dataclass
class _FakeResponse:
    content: list[_FakeBlock]


class FakeAnthropicClient:
    """Returns scripted JSON responses captured by the test."""

    def __init__(self, scripted_response: dict):
        self._scripted = scripted_response
        self.last_messages: list[dict] | None = None
        self.last_system: object | None = None

    @property
    def messages(self):
        return self  # so client.messages.create works

    def create(self, **kwargs):
        self.last_messages = kwargs.get("messages")
        self.last_system = kwargs.get("system")
        text = json.dumps(self._scripted)
        return _FakeResponse(content=[_FakeBlock(type="text", text=text)])


@pytest.fixture
def fake_client_factory():
    created: list[FakeAnthropicClient] = []

    def make(decisions: list[dict]):
        client = FakeAnthropicClient({"decisions": decisions})
        created.append(client)
        set_ai_client_override(client)
        return client

    yield make
    set_ai_client_override(None)


# ── Helper builders ───────────────────────────────────────────────────


def _make_deck(slides: list[list[str]]) -> bytes:
    """Build a deck where each slide is a list of text-frame strings.

    First string on each slide acts as the title.
    """
    prs = Presentation()
    for slide_texts in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        top = Inches(0.5)
        for text in slide_texts:
            box = slide.shapes.add_textbox(Inches(0.5), top, Inches(7), Inches(0.6))
            run = box.text_frame.paragraphs[0].add_run()
            run.text = text
            top += Inches(0.7)
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _deck_text(file_bytes: bytes) -> str:
    prs = Presentation(BytesIO(file_bytes))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return " ".join(parts)


def _create_period(db, *, status: str = "draft") -> ReportingPeriod:
    upload = CsvUpload(
        filename="phase_b.csv",
        status="complete",
        stichtag=date(2025, 6, 30),
        row_count=1,
        data_row_count=1,
        summary_row_count=0,
    )
    db.add(upload)
    db.commit()
    db.refresh(upload)
    db.add(
        RawRentRoll(
            upload_id=upload.id,
            row_number=1,
            row_type="data",
            fund="GLIF",
            property_id="2001",
            tenant_name="Tenant Z",
            unit_type="Halle",
            area_sqm=2000,
            annual_net_rent=240000,
            monthly_net_rent=20000,
        )
    )
    period = ReportingPeriod(stichtag=date(2025, 6, 30), upload_id=upload.id, status=status)
    db.add(period)
    db.commit()
    db.refresh(period)
    return period


def _wait_for(client, job_id: int, status: str = "proposed", timeout: float = 4.0) -> dict:
    deadline = time.time() + timeout
    last = None
    while time.time() < deadline:
        response = client.get(f"/api/pptx/{job_id}")
        assert response.status_code == 200
        last = response.json()
        if last["status"] in (status, "error"):
            return last
        time.sleep(0.05)
    raise AssertionError(f"job {job_id} stuck at {last and last['status']}")


# ── Candidate extraction (unit) ───────────────────────────────────────


def test_collect_candidates_filters_dates_and_titles():
    deck = _make_deck([["Portfolio overview Q2 2025", "Total rent", "12,5 M€", "30.06.2025"]])
    elements = ingest_pptx(deck)
    candidates = collect_candidates(elements)
    originals = [c.original_value for c in candidates]
    assert "12,5 M€" in originals
    assert "30.06.2025" not in originals  # date filtered out
    assert "Total rent" not in originals
    rent_candidate = next(c for c in candidates if c.original_value == "12,5 M€")
    assert "Total rent" in rent_candidate.label_context
    assert "Portfolio overview" in rent_candidate.slide_title


# ── resolve_with_ai (unit, no API) ────────────────────────────────────


def test_resolve_with_ai_mapping(db, fake_client_factory):
    period = _create_period(db)
    deck = _make_deck([["Q2 portfolio overview", "Total rent", "10,0 M€"]])
    elements = ingest_pptx(deck)
    client = fake_client_factory(
        [
            {
                "idx": 0,
                "decision": "mapping",
                "kpi_id": "total_rent",
                "confidence": 0.95,
                "reasoning": "Adjacent label says 'Total rent'.",
            }
        ]
    )

    result = resolve_with_ai(db, elements, period.id, client=client)

    assert result["mode"] == "ai"
    assert result["summary"]["mappings"] == 1
    proposal = result["proposals"][0]
    assert proposal["kind"] == "mapping"
    assert proposal["kpi_id"] == "total_rent"
    assert proposal["new_value"] == "0,2 M€"  # 240k EUR formatted
    assert "total_rent" in result["available_kpis"]


def test_resolve_with_ai_ambiguous_scope(db, fake_client_factory):
    period = _create_period(db)
    deck = _make_deck([["Rent overview", "Rent", "5,0 M€"]])
    elements = ingest_pptx(deck)
    client = fake_client_factory(
        [
            {
                "idx": 0,
                "decision": "ambiguous_scope",
                "kpi_id": "total_rent",
                "confidence": 0.4,
                "reasoning": "Cannot tell whether portfolio or fund.",
            }
        ]
    )

    result = resolve_with_ai(db, elements, period.id, client=client)
    proposal = result["proposals"][0]
    assert proposal["kind"] == "ambiguous_scope"
    assert proposal["candidate_kpi_id"] == "total_rent"
    assert "new_value" not in proposal  # no value committed for ambiguous
    assert result["summary"]["ambiguous"] == 1


def test_resolve_with_ai_unsupported(db, fake_client_factory):
    period = _create_period(db)
    deck = _make_deck([["Performance", "NOI", "2,1 M€"]])
    elements = ingest_pptx(deck)
    client = fake_client_factory(
        [
            {
                "idx": 0,
                "decision": "unsupported_kpi",
                "label_observed": "NOI",
                "confidence": 0.9,
            }
        ]
    )

    result = resolve_with_ai(db, elements, period.id, client=client)
    proposal = result["proposals"][0]
    assert proposal["kind"] == "unsupported_kpi"
    assert proposal["label_observed"] == "NOI"
    assert result["summary"]["unsupported"] == 1


def test_resolve_with_ai_mapping_falls_back_when_value_missing(db, fake_client_factory):
    """If the AI picks a KPI that has no value for the period, downgrade to unsupported."""
    period = _create_period(db)
    deck = _make_deck([["Performance", "Fair value", "100 M€"]])
    elements = ingest_pptx(deck)
    client = fake_client_factory(
        [
            {
                "idx": 0,
                "decision": "mapping",
                "kpi_id": "fair_value",
                "confidence": 0.99,
            }
        ]
    )

    result = resolve_with_ai(db, elements, period.id, client=client)
    proposal = result["proposals"][0]
    # Draft period without snapshot data => fair_value is None => downgrade.
    assert proposal["kind"] == "unsupported_kpi"


def test_resolve_with_ai_missing_decision_marked_skipped(db, fake_client_factory):
    period = _create_period(db)
    deck = _make_deck([["Overview", "Total rent", "10 M€", "Total area", "1.500 m²"]])
    elements = ingest_pptx(deck)
    client = fake_client_factory(
        [
            {"idx": 0, "decision": "mapping", "kpi_id": "total_rent", "confidence": 0.9},
        ]
    )

    result = resolve_with_ai(db, elements, period.id, client=client)
    kinds = [p["kind"] for p in result["proposals"]]
    assert kinds[0] == "mapping"
    # Second candidate has no decision returned by the AI -> skipped.
    assert kinds[1] == "skipped"


# ── End-to-end API (scan + apply) ─────────────────────────────────────


def test_scan_endpoint_runs_ai_and_apply_works(client, db, pptx_upload_dir, fake_client_factory):
    period = _create_period(db, status="finalized")
    deck = _make_deck([["Q2 overview", "Total rent", "12,5 M€"]])

    upload = client.post(
        "/api/pptx/upload",
        files={"file": ("phaseb.pptx", deck, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    assert upload.status_code == 200
    job_id = upload.json()["id"]
    _wait_for(client, job_id, status="proposed")

    fake_client_factory(
        [
            {
                "idx": 0,
                "decision": "mapping",
                "kpi_id": "total_rent",
                "confidence": 0.96,
                "reasoning": "Adjacent label 'Total rent'.",
            }
        ]
    )

    scan = client.post(f"/api/pptx/{job_id}/scan", params={"period_id": period.id})
    assert scan.status_code == 200
    job = _wait_for(client, job_id, status="proposed")
    assert job["proposals"]["mode"] == "ai"
    assert job["proposals"]["summary"]["mappings"] == 1
    proposal = job["proposals"]["ai_proposals"][0]
    assert proposal["kpi_id"] == "total_rent"
    assert proposal["new_value"]

    apply_resp = client.post(
        f"/api/pptx/{job_id}/apply",
        json={
            "period_id": period.id,
            "ai_confirmations": [{"idx": proposal["idx"]}],
        },
    )
    assert apply_resp.status_code == 200
    data = apply_resp.json()
    assert data["status"] == "complete"
    assert data["confirmed_json"][0]["mode"] == "ai"
    assert data["confirmed_json"][0]["kpi_id"] == "total_rent"

    download = client.get(f"/api/pptx/{job_id}/download")
    assert download.status_code == 200
    assert proposal["new_value"] in _deck_text(download.content)
    assert "12,5 M€" not in _deck_text(download.content)


def test_apply_ai_with_ambiguous_scope_skip_and_portfolio(client, db, pptx_upload_dir, fake_client_factory):
    period = _create_period(db, status="finalized")
    deck = _make_deck(
        [
            ["Overview", "Rent", "5,0 M€", "Vacancy rate", "8 %"],
        ]
    )

    upload = client.post(
        "/api/pptx/upload",
        files={"file": ("ambig.pptx", deck, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    job_id = upload.json()["id"]
    _wait_for(client, job_id, status="proposed")

    fake_client_factory(
        [
            {
                "idx": 0,
                "decision": "ambiguous_scope",
                "kpi_id": "total_rent",
                "confidence": 0.4,
            },
            {
                "idx": 1,
                "decision": "mapping",
                "kpi_id": "vacancy_rate",
                "confidence": 0.9,
            },
        ]
    )

    scan = client.post(f"/api/pptx/{job_id}/scan", params={"period_id": period.id})
    assert scan.status_code == 200
    job = _wait_for(client, job_id, status="proposed")
    proposals = job["proposals"]["ai_proposals"]
    assert {p["kind"] for p in proposals} == {"ambiguous_scope", "mapping"}

    # Reject ambiguous proposal explicitly via scope_choice=skip; accept mapping.
    apply_resp = client.post(
        f"/api/pptx/{job_id}/apply",
        json={
            "period_id": period.id,
            "ai_confirmations": [
                {"idx": 0, "scope_choice": "skip"},
                {"idx": 1},
            ],
        },
    )
    assert apply_resp.status_code == 200
    data = apply_resp.json()
    assert data["status"] == "complete"
    confirmed_idxs = [c["idx"] for c in data["confirmed_json"]]
    assert confirmed_idxs == [1]


def test_apply_ai_ambiguous_requires_kpi(client, db, pptx_upload_dir, fake_client_factory):
    period = _create_period(db, status="finalized")
    deck = _make_deck([["Overview", "Total", "5,0 M€"]])

    upload = client.post(
        "/api/pptx/upload",
        files={"file": ("ambig2.pptx", deck, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    job_id = upload.json()["id"]
    _wait_for(client, job_id, status="proposed")

    fake_client_factory(
        [
            {
                "idx": 0,
                "decision": "ambiguous_scope",
                "kpi_id": None,
                "confidence": 0.3,
            }
        ]
    )
    client.post(f"/api/pptx/{job_id}/scan", params={"period_id": period.id})
    _wait_for(client, job_id, status="proposed")

    # scope_choice=portfolio without a kpi_id -> error (fails closed).
    bad = client.post(
        f"/api/pptx/{job_id}/apply",
        json={
            "period_id": period.id,
            "ai_confirmations": [{"idx": 0, "scope_choice": "portfolio"}],
        },
    )
    assert bad.status_code == 200
    assert bad.json()["status"] == "error"
    assert "kpi_id" in bad.json()["error_message"]


def test_apply_ai_no_confirmations_errors(client, db, pptx_upload_dir, fake_client_factory):
    period = _create_period(db, status="finalized")
    deck = _make_deck([["Overview", "Total rent", "5,0 M€"]])
    upload = client.post(
        "/api/pptx/upload",
        files={"file": ("none.pptx", deck, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    job_id = upload.json()["id"]
    _wait_for(client, job_id, status="proposed")

    fake_client_factory(
        [
            {
                "idx": 0,
                "decision": "unsupported_kpi",
                "label_observed": "Total rent",
            }
        ]
    )
    client.post(f"/api/pptx/{job_id}/scan", params={"period_id": period.id})
    _wait_for(client, job_id, status="proposed")

    # Empty confirmations list -> fails closed.
    resp = client.post(
        f"/api/pptx/{job_id}/apply",
        json={"period_id": period.id, "ai_confirmations": []},
    )
    assert resp.status_code == 200
    assert resp.json()["status"] == "error"


def test_apply_ai_ambiguous_rejects_implicit_candidate_kpi(
    client, db, pptx_upload_dir, fake_client_factory
):
    """Codex P2: even when the AI returned `candidate_kpi_id`, the user must explicitly confirm."""
    period = _create_period(db, status="finalized")
    deck = _make_deck([["Overview", "Total", "5,0 M€"]])

    upload = client.post(
        "/api/pptx/upload",
        files={"file": ("ambig3.pptx", deck, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    job_id = upload.json()["id"]
    _wait_for(client, job_id, status="proposed")

    fake_client_factory(
        [
            {
                "idx": 0,
                "decision": "ambiguous_scope",
                "kpi_id": "total_rent",  # AI suggests a candidate
                "confidence": 0.4,
            }
        ]
    )
    client.post(f"/api/pptx/{job_id}/scan", params={"period_id": period.id})
    job = _wait_for(client, job_id, status="proposed")
    assert job["proposals"]["ai_proposals"][0]["candidate_kpi_id"] == "total_rent"

    # Sending only scope_choice=portfolio without kpi_id must fail closed.
    bad = client.post(
        f"/api/pptx/{job_id}/apply",
        json={
            "period_id": period.id,
            "ai_confirmations": [{"idx": 0, "scope_choice": "portfolio"}],
        },
    )
    assert bad.status_code == 200
    assert bad.json()["status"] == "error"
    assert "kpi_id" in bad.json()["error_message"]


def test_apply_ai_re_resolves_for_apply_period(client, db, pptx_upload_dir, fake_client_factory):
    """Codex P2: changing the period between scan and apply must reflect new values."""
    # Two finalized periods with very different rent totals.
    period_a = _create_period(db, status="finalized")
    upload_b = CsvUpload(
        filename="period_b.csv",
        status="complete",
        stichtag=date(2025, 9, 30),
        row_count=1,
        data_row_count=1,
        summary_row_count=0,
    )
    db.add(upload_b)
    db.commit()
    db.refresh(upload_b)
    db.add(
        RawRentRoll(
            upload_id=upload_b.id,
            row_number=1,
            row_type="data",
            fund="GLIF",
            property_id="3001",
            tenant_name="Tenant Q",
            unit_type="Halle",
            area_sqm=4000,
            annual_net_rent=999_000_000,
            monthly_net_rent=83_250_000,
        )
    )
    period_b = ReportingPeriod(
        stichtag=date(2025, 9, 30), upload_id=upload_b.id, status="finalized"
    )
    db.add(period_b)
    db.commit()
    db.refresh(period_b)

    deck = _make_deck([["Overview", "Total rent", "1,0 M€"]])
    upload = client.post(
        "/api/pptx/upload",
        files={"file": ("reapply.pptx", deck, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    job_id = upload.json()["id"]
    _wait_for(client, job_id, status="proposed")

    fake_client_factory(
        [{"idx": 0, "decision": "mapping", "kpi_id": "total_rent", "confidence": 0.95}]
    )
    # Scan against period A.
    client.post(f"/api/pptx/{job_id}/scan", params={"period_id": period_a.id})
    job = _wait_for(client, job_id, status="proposed")
    proposal = job["proposals"]["ai_proposals"][0]
    new_value_a = proposal["new_value"]

    # Apply against period B - the value must be re-resolved, not the cached one.
    apply_resp = client.post(
        f"/api/pptx/{job_id}/apply",
        json={
            "period_id": period_b.id,
            "ai_confirmations": [{"idx": 0}],
        },
    )
    assert apply_resp.status_code == 200
    data = apply_resp.json()
    assert data["status"] == "complete"
    assert data["reporting_period_id"] == period_b.id
    confirmed = data["confirmed_json"][0]
    # Period B's annual rent is 999M EUR (vs ~0,2 M€ for period A); the patcher
    # must re-resolve to that fresh value rather than the stored proposal one.
    assert confirmed["new_value"] != new_value_a
    assert confirmed["new_value"] == "999 M€"

    download = client.get(f"/api/pptx/{job_id}/download")
    assert "999 M€" in _deck_text(download.content)


def test_collect_candidates_currency_prefix():
    """Codex P2: ``EUR 12,500,000`` and ``€ 12,5 M`` should be candidates."""
    deck = _make_deck(
        [
            [
                "Overview",
                "Total rent",
                "EUR 12,500,000",
                "Fair value",
                "€ 250 M",
            ]
        ]
    )
    elements = ingest_pptx(deck)
    candidates = collect_candidates(elements)
    originals = [c.original_value for c in candidates]
    assert "EUR 12,500,000" in originals
    assert "€ 250 M" in originals


def test_scan_records_period_status_for_draft(client, db, pptx_upload_dir, fake_client_factory):
    period = _create_period(db, status="draft")
    deck = _make_deck([["Overview", "Total rent", "5,0 M€"]])

    upload = client.post(
        "/api/pptx/upload",
        files={"file": ("draft.pptx", deck, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    job_id = upload.json()["id"]
    _wait_for(client, job_id, status="proposed")

    fake_client_factory(
        [
            {
                "idx": 0,
                "decision": "mapping",
                "kpi_id": "total_rent",
                "confidence": 0.9,
            }
        ]
    )
    client.post(f"/api/pptx/{job_id}/scan", params={"period_id": period.id})
    job = _wait_for(client, job_id, status="proposed")
    assert job["proposals"]["mode"] == "ai"

    apply_resp = client.post(
        f"/api/pptx/{job_id}/apply",
        json={
            "period_id": period.id,
            "ai_confirmations": [{"idx": 0}],
        },
    )
    assert apply_resp.json()["period_status_at_refresh"] == "draft"
