"""PPTX slide generation for property factsheets, fund summaries, and portfolio reports."""
from __future__ import annotations

import io
from datetime import date

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from sqlalchemy import func
from sqlalchemy.orm import Session

from app.models.database import (
    CsvUpload,
    PropertyMaster,
    RawRentRoll,
    ReportingPeriod,
    SnapshotPropertyMaster,
)

GARBE_BLAU = RGBColor(0x00, 0x32, 0x55)
GARBE_GRUN = RGBColor(0x64, 0xB4, 0x2D)
GARBE_OCKER = RGBColor(0xA4, 0x81, 0x13)
GARBE_ROT = RGBColor(0xFF, 0x72, 0x76)
GARBE_TURKIS = RGBColor(0x00, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEC, 0xEC, 0xEC)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _new_pres() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def _add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = GARBE_BLAU

    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xC0, 0xCA, 0xDA)
        p2.alignment = PP_ALIGN.LEFT

    return slide


def _add_content_slide(prs: Presentation, title: str) -> tuple:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header_shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = GARBE_BLAU
    header_shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE

    return slide


def _add_kpi_box(slide, left, top, width, height, label: str, value: str, color=GARBE_BLAU):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LIGHT_GRAY
    shape.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), Inches(0.3),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = label.upper()
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x87, 0x9C, 0xB5)
    p.font.bold = True

    txBox2 = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.45),
        width - Inches(0.3), Inches(0.5),
    )
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = value
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = color


def _fmt_eur(v) -> str:
    if v is None:
        return "—"
    return f"€{v:,.0f}".replace(",", ".")


def _fmt_num(v, decimals=0) -> str:
    if v is None:
        return "—"
    if decimals == 0:
        return f"{v:,.0f}".replace(",", ".")
    return f"{v:,.{decimals}f}".replace(",", "X").replace(".", ",").replace("X", ".")


def _fmt_pct(v) -> str:
    if v is None:
        return "—"
    return f"{v:.1f}%"


def _add_table(slide, left, top, width, headers: list[str], rows: list[list[str]], col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.3 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = GARBE_BLAU
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF9, 0xF9, 0xF9)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = GARBE_BLAU

    return table_shape


def _chart_to_image(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, transparent=False)
    plt.close(fig)
    buf.seek(0)
    return buf


def _add_bar_chart(slide, left, top, width, height, categories: list[str], values: list[float],
                   title: str = "", series_name: str = "Value", chart_type=XL_CHART_TYPE.BAR_CLUSTERED):
    """Add a native pptx bar chart to a slide."""
    from pptx.chart.data import CategoryChartData
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)

    chart_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = chart_frame.chart
    chart.has_legend = False

    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = GARBE_BLAU

    plot = chart.plots[0]
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = GARBE_BLAU

    return chart


def _get_property_data(db: Session, upload_id: int, property_id: str) -> dict:
    """Get aggregated data for a single property from an upload."""
    base = db.query(RawRentRoll).filter(
        RawRentRoll.upload_id == upload_id,
        RawRentRoll.row_type.in_(["data", "orphan"]),
        RawRentRoll.property_id == property_id,
    )

    total_rent = float(base.with_entities(func.coalesce(func.sum(RawRentRoll.annual_net_rent), 0)).scalar() or 0)
    total_area = float(base.with_entities(func.coalesce(func.sum(RawRentRoll.area_sqm), 0)).scalar() or 0)
    vacant_area = float(
        base.filter(RawRentRoll.tenant_name == "LEERSTAND")
        .with_entities(func.coalesce(func.sum(RawRentRoll.area_sqm), 0)).scalar() or 0
    )
    tenant_count = base.filter(
        RawRentRoll.tenant_name.isnot(None),
        RawRentRoll.tenant_name != "LEERSTAND",
    ).with_entities(func.count(func.distinct(RawRentRoll.tenant_name))).scalar() or 0

    area_by_type = {}
    type_rows = (
        base.with_entities(RawRentRoll.unit_type, func.sum(RawRentRoll.area_sqm))
        .filter(RawRentRoll.unit_type.isnot(None))
        .group_by(RawRentRoll.unit_type)
        .all()
    )
    for ut, area in type_rows:
        area_by_type[ut] = float(area or 0)

    top_tenants = (
        base.filter(
            RawRentRoll.tenant_name.isnot(None),
            RawRentRoll.tenant_name != "LEERSTAND",
        )
        .with_entities(
            RawRentRoll.tenant_name,
            func.sum(RawRentRoll.annual_net_rent),
            func.sum(RawRentRoll.area_sqm),
        )
        .group_by(RawRentRoll.tenant_name)
        .order_by(func.sum(RawRentRoll.annual_net_rent).desc())
        .limit(10)
        .all()
    )

    wault = (
        db.query(RawRentRoll.wault)
        .filter(
            RawRentRoll.upload_id == upload_id,
            RawRentRoll.row_type == "property_summary",
            RawRentRoll.property_id == property_id,
            RawRentRoll.wault.isnot(None),
        )
        .scalar()
    )

    prop_master = db.query(PropertyMaster).filter(PropertyMaster.property_id == property_id).first()

    return {
        "total_rent": total_rent,
        "total_area": total_area,
        "vacant_area": vacant_area,
        "vacancy_rate": (vacant_area / total_area * 100) if total_area > 0 else 0,
        "tenant_count": tenant_count,
        "area_by_type": area_by_type,
        "top_tenants": [(t, float(r or 0), float(a or 0)) for t, r, a in top_tenants],
        "wault": float(wault) if wault else None,
        "city": prop_master.city if prop_master else None,
        "country": prop_master.country if prop_master else None,
        "fair_value": float(prop_master.fair_value) if prop_master and prop_master.fair_value else None,
        "construction_year": prop_master.construction_year if prop_master else None,
    }


def generate_property_factsheet(db: Session, upload_id: int, property_id: str) -> io.BytesIO:
    """Generate a 2-slide PPTX factsheet for a single property."""
    upload = db.get(CsvUpload, upload_id)
    data = _get_property_data(db, upload_id, property_id)
    stichtag = upload.stichtag if upload else None

    prs = _new_pres()

    location = data["city"] or ""
    if data["country"]:
        location = f"{location}, {data['country']}" if location else data["country"]

    _add_title_slide(
        prs,
        f"Property {property_id}",
        f"{location} — {stichtag or 'N/A'}",
    )

    # Slide 2: KPIs + area breakdown
    slide = _add_content_slide(prs, f"Property {property_id} — Key Metrics")

    kpi_y = Inches(1.2)
    kpi_w = Inches(2.8)
    kpi_h = Inches(1.0)
    gap = Inches(0.2)

    _add_kpi_box(slide, Inches(0.5), kpi_y, kpi_w, kpi_h, "Annual Rent", _fmt_eur(data["total_rent"]))
    _add_kpi_box(slide, Inches(0.5) + kpi_w + gap, kpi_y, kpi_w, kpi_h, "Rentable Area", f"{_fmt_num(data['total_area'])} sqm")
    _add_kpi_box(slide, Inches(0.5) + 2 * (kpi_w + gap), kpi_y, kpi_w, kpi_h, "Vacancy Rate", _fmt_pct(data["vacancy_rate"]),
                 color=GARBE_GRUN if data["vacancy_rate"] < 5 else GARBE_OCKER if data["vacancy_rate"] < 15 else GARBE_ROT)
    _add_kpi_box(slide, Inches(0.5) + 3 * (kpi_w + gap), kpi_y, kpi_w, kpi_h, "WAULT", f"{_fmt_num(data['wault'], 1)} yrs" if data["wault"] else "—")

    row2_y = kpi_y + kpi_h + gap
    _add_kpi_box(slide, Inches(0.5), row2_y, kpi_w, kpi_h, "Tenants", str(data["tenant_count"]))
    _add_kpi_box(slide, Inches(0.5) + kpi_w + gap, row2_y, kpi_w, kpi_h, "Fair Value", _fmt_eur(data["fair_value"]))
    if data["construction_year"]:
        _add_kpi_box(slide, Inches(0.5) + 2 * (kpi_w + gap), row2_y, kpi_w, kpi_h, "Built", str(data["construction_year"]))

    # Area breakdown chart
    if data["area_by_type"]:
        fig, ax = plt.subplots(figsize=(4, 3))
        colors = ["#003255", "#64B42D", "#005555", "#A48113", "#FF7276", "#337777", "#879CB5", "#99BF65", "#669999", "#C0CADA"]
        types = list(data["area_by_type"].keys())
        areas = list(data["area_by_type"].values())
        wedges, texts, autotexts = ax.pie(
            areas, labels=types, autopct="%1.0f%%",
            colors=colors[:len(types)],
            textprops={"fontsize": 8},
        )
        ax.set_title("Area by Type", fontsize=10, fontweight="bold", color="#003255")
        img_buf = _chart_to_image(fig)
        slide.shapes.add_picture(img_buf, Inches(0.5), row2_y + kpi_h + gap, Inches(5), Inches(3.5))

    # Top tenants table
    if data["top_tenants"]:
        slide2 = _add_content_slide(prs, f"Property {property_id} — Top Tenants")
        headers = ["Tenant", "Annual Rent", "Area (sqm)", "% of Rent"]
        rows = []
        for name, rent, area in data["top_tenants"]:
            pct = (rent / data["total_rent"] * 100) if data["total_rent"] > 0 else 0
            rows.append([name[:40], _fmt_eur(rent), _fmt_num(area), _fmt_pct(pct)])
        _add_table(slide2, Inches(0.5), Inches(1.2), Inches(12), headers, rows)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _get_portfolio_data(db: Session, upload_id: int) -> dict:
    """Get portfolio-level aggregated data."""
    base = db.query(RawRentRoll).filter(
        RawRentRoll.upload_id == upload_id,
        RawRentRoll.row_type.in_(["data", "orphan"]),
    )

    total_rent = float(base.with_entities(func.coalesce(func.sum(RawRentRoll.annual_net_rent), 0)).scalar() or 0)
    total_area = float(base.with_entities(func.coalesce(func.sum(RawRentRoll.area_sqm), 0)).scalar() or 0)
    vacant_area = float(
        base.filter(RawRentRoll.tenant_name == "LEERSTAND")
        .with_entities(func.coalesce(func.sum(RawRentRoll.area_sqm), 0)).scalar() or 0
    )
    tenant_count = base.filter(
        RawRentRoll.tenant_name.isnot(None),
        RawRentRoll.tenant_name != "LEERSTAND",
    ).with_entities(func.count(func.distinct(RawRentRoll.tenant_name))).scalar() or 0
    property_count = base.with_entities(func.count(func.distinct(RawRentRoll.property_id))).scalar() or 0

    properties = (
        base.with_entities(
            RawRentRoll.property_id,
            func.sum(RawRentRoll.annual_net_rent),
            func.sum(RawRentRoll.area_sqm),
        )
        .group_by(RawRentRoll.property_id)
        .order_by(func.sum(RawRentRoll.annual_net_rent).desc())
        .all()
    )

    rent_by_fund = (
        base.with_entities(RawRentRoll.fund, func.sum(RawRentRoll.annual_net_rent))
        .group_by(RawRentRoll.fund)
        .order_by(func.sum(RawRentRoll.annual_net_rent).desc())
        .all()
    )

    top_tenants = (
        base.filter(
            RawRentRoll.tenant_name.isnot(None),
            RawRentRoll.tenant_name != "LEERSTAND",
        )
        .with_entities(
            RawRentRoll.tenant_name,
            func.sum(RawRentRoll.annual_net_rent),
        )
        .group_by(RawRentRoll.tenant_name)
        .order_by(func.sum(RawRentRoll.annual_net_rent).desc())
        .limit(10)
        .all()
    )

    return {
        "total_rent": total_rent,
        "total_area": total_area,
        "vacant_area": vacant_area,
        "vacancy_rate": (vacant_area / total_area * 100) if total_area > 0 else 0,
        "tenant_count": tenant_count,
        "property_count": property_count,
        "properties": [(pid, float(r or 0), float(a or 0)) for pid, r, a in properties],
        "rent_by_fund": [(f, float(r or 0)) for f, r in rent_by_fund],
        "top_tenants": [(t, float(r or 0)) for t, r in top_tenants],
    }


def generate_portfolio_overview(db: Session, upload_id: int) -> io.BytesIO:
    """Generate a multi-slide portfolio overview PPTX."""
    upload = db.get(CsvUpload, upload_id)
    data = _get_portfolio_data(db, upload_id)
    stichtag = upload.stichtag if upload else None

    prs = _new_pres()
    _add_title_slide(prs, "Portfolio Overview", f"Stichtag: {stichtag or 'N/A'}")

    # KPI slide
    slide = _add_content_slide(prs, "Portfolio KPIs")
    kpi_y = Inches(1.2)
    kpi_w = Inches(2.8)
    kpi_h = Inches(1.0)
    gap = Inches(0.2)

    _add_kpi_box(slide, Inches(0.5), kpi_y, kpi_w, kpi_h, "Total Annual Rent", _fmt_eur(data["total_rent"]))
    _add_kpi_box(slide, Inches(0.5) + kpi_w + gap, kpi_y, kpi_w, kpi_h, "Total Area", f"{_fmt_num(data['total_area'])} sqm")
    _add_kpi_box(slide, Inches(0.5) + 2 * (kpi_w + gap), kpi_y, kpi_w, kpi_h, "Vacancy Rate", _fmt_pct(data["vacancy_rate"]),
                 color=GARBE_GRUN if data["vacancy_rate"] < 5 else GARBE_OCKER)
    _add_kpi_box(slide, Inches(0.5) + 3 * (kpi_w + gap), kpi_y, kpi_w, kpi_h, "Properties", str(data["property_count"]))

    row2_y = kpi_y + kpi_h + gap
    _add_kpi_box(slide, Inches(0.5), row2_y, kpi_w, kpi_h, "Tenants", str(data["tenant_count"]))

    # Rent by fund chart
    if data["rent_by_fund"]:
        funds = [f[:20] for f, _ in data["rent_by_fund"]]
        rents = [r / 1000 for _, r in data["rent_by_fund"]]
        _add_bar_chart(slide, Inches(6.5), row2_y, Inches(6), Inches(3.5),
                       funds, rents, title="Rent by Fund (k€)", series_name="Rent",
                       chart_type=XL_CHART_TYPE.BAR_CLUSTERED)

    # Top tenants slide
    if data["top_tenants"]:
        slide2 = _add_content_slide(prs, "Top 10 Tenants by Rent")
        names = [t[:30] for t, _ in data["top_tenants"]]
        rents = [r / 1000 for _, r in data["top_tenants"]]
        _add_bar_chart(slide2, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5),
                       names, rents, title="Annual Rent (k€)", series_name="Rent",
                       chart_type=XL_CHART_TYPE.BAR_CLUSTERED)

    # Property table slide
    if data["properties"]:
        slide3 = _add_content_slide(prs, "Properties by Rent")
        headers = ["Property ID", "Annual Rent", "Area (sqm)", "% of Total"]
        rows = []
        for pid, rent, area in data["properties"][:20]:
            pct = (rent / data["total_rent"] * 100) if data["total_rent"] > 0 else 0
            pm = db.query(PropertyMaster).filter(PropertyMaster.property_id == pid).first()
            label = f"{pid} — {pm.city}" if pm and pm.city else str(pid)
            rows.append([label, _fmt_eur(rent), _fmt_num(area), _fmt_pct(pct)])
        _add_table(slide3, Inches(0.5), Inches(1.2), Inches(12), headers, rows)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def generate_lease_expiry_profile(db: Session, upload_id: int) -> io.BytesIO:
    """Generate a lease expiry waterfall chart PPTX."""
    upload = db.get(CsvUpload, upload_id)
    stichtag = upload.stichtag if upload else date.today()

    data_rows = (
        db.query(RawRentRoll)
        .filter(
            RawRentRoll.upload_id == upload_id,
            RawRentRoll.row_type.in_(["data", "orphan"]),
            RawRentRoll.tenant_name != "LEERSTAND",
            RawRentRoll.annual_net_rent.isnot(None),
        )
        .all()
    )

    buckets = {i: 0.0 for i in range(11)}
    open_ended = 0.0

    for row in data_rows:
        rent = float(row.annual_net_rent or 0)
        if rent <= 0:
            continue
        lease_end = row.lease_end_agreed or row.lease_end_actual
        if not lease_end:
            open_ended += rent
            continue
        years_to_expiry = lease_end.year - stichtag.year
        if years_to_expiry < 0:
            buckets[0] += rent
        elif years_to_expiry >= 10:
            buckets[10] += rent
        else:
            buckets[years_to_expiry] += rent

    prs = _new_pres()
    _add_title_slide(prs, "Lease Expiry Profile", f"Stichtag: {stichtag}")

    slide = _add_content_slide(prs, "Lease Expiry by Year")

    labels = [f"Year {i}" for i in range(10)] + ["10+", "Open"]
    values = [buckets[i] / 1000 for i in range(11)] + [open_ended / 1000]

    _add_bar_chart(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5),
                   labels, values, title="Annual Rent by Lease Expiry Year (k€)",
                   series_name="Rent Expiring",
                   chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED)

    # Summary table
    total_rent = sum(buckets.values()) + open_ended
    slide2 = _add_content_slide(prs, "Lease Expiry Summary")
    headers = ["Bucket", "Rent Expiring", "% of Total"]
    rows = []
    for i in range(11):
        label = f"Year {i}" if i < 10 else "10+ years"
        pct = (buckets[i] / total_rent * 100) if total_rent > 0 else 0
        rows.append([label, _fmt_eur(buckets[i]), _fmt_pct(pct)])
    if open_ended > 0:
        pct = (open_ended / total_rent * 100) if total_rent > 0 else 0
        rows.append(["Open-ended", _fmt_eur(open_ended), _fmt_pct(pct)])
    rows.append(["Total", _fmt_eur(total_rent), "100.0%"])
    _add_table(slide2, Inches(0.5), Inches(1.2), Inches(8), headers, rows)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def generate_fund_summary(db: Session, upload_id: int, fund_name: str) -> io.BytesIO:
    """Generate a fund-level summary PPTX."""
    upload = db.get(CsvUpload, upload_id)
    stichtag = upload.stichtag if upload else None

    base = db.query(RawRentRoll).filter(
        RawRentRoll.upload_id == upload_id,
        RawRentRoll.row_type.in_(["data", "orphan"]),
        RawRentRoll.fund == fund_name,
    )

    total_rent = float(base.with_entities(func.coalesce(func.sum(RawRentRoll.annual_net_rent), 0)).scalar() or 0)
    total_area = float(base.with_entities(func.coalesce(func.sum(RawRentRoll.area_sqm), 0)).scalar() or 0)
    vacant_area = float(
        base.filter(RawRentRoll.tenant_name == "LEERSTAND")
        .with_entities(func.coalesce(func.sum(RawRentRoll.area_sqm), 0)).scalar() or 0
    )
    tenant_count = base.filter(
        RawRentRoll.tenant_name.isnot(None),
        RawRentRoll.tenant_name != "LEERSTAND",
    ).with_entities(func.count(func.distinct(RawRentRoll.tenant_name))).scalar() or 0
    property_count = base.with_entities(func.count(func.distinct(RawRentRoll.property_id))).scalar() or 0

    properties = (
        base.with_entities(
            RawRentRoll.property_id,
            func.sum(RawRentRoll.annual_net_rent),
            func.sum(RawRentRoll.area_sqm),
        )
        .group_by(RawRentRoll.property_id)
        .order_by(func.sum(RawRentRoll.annual_net_rent).desc())
        .all()
    )

    prs = _new_pres()
    _add_title_slide(prs, f"Fund: {fund_name}", f"Stichtag: {stichtag or 'N/A'}")

    # KPIs
    slide = _add_content_slide(prs, f"{fund_name} — Key Metrics")
    kpi_y = Inches(1.2)
    kpi_w = Inches(2.8)
    kpi_h = Inches(1.0)
    gap = Inches(0.2)

    _add_kpi_box(slide, Inches(0.5), kpi_y, kpi_w, kpi_h, "Annual Rent", _fmt_eur(total_rent))
    _add_kpi_box(slide, Inches(0.5) + kpi_w + gap, kpi_y, kpi_w, kpi_h, "Total Area", f"{_fmt_num(total_area)} sqm")
    _add_kpi_box(slide, Inches(0.5) + 2 * (kpi_w + gap), kpi_y, kpi_w, kpi_h, "Vacancy",
                 _fmt_pct(vacant_area / total_area * 100 if total_area > 0 else 0))
    _add_kpi_box(slide, Inches(0.5) + 3 * (kpi_w + gap), kpi_y, kpi_w, kpi_h, "Properties", str(property_count))

    # Property table
    if properties:
        slide2 = _add_content_slide(prs, f"{fund_name} — Properties")
        headers = ["Property ID", "City", "Annual Rent", "Area (sqm)"]
        rows = []
        for pid, rent, area in properties:
            pm = db.query(PropertyMaster).filter(PropertyMaster.property_id == pid).first()
            city = pm.city if pm else "—"
            rows.append([str(pid), city, _fmt_eur(float(rent or 0)), _fmt_num(float(area or 0))])
        _add_table(slide2, Inches(0.5), Inches(1.2), Inches(12), headers, rows[:20])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
