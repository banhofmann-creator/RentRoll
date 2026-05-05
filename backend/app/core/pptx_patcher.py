from __future__ import annotations

import io
from dataclasses import asdict, dataclass, is_dataclass
from typing import Any

from pptx import Presentation

from app.parsers.pptx_ingestor import TextAddress


@dataclass(frozen=True)
class Mapping:
    address: TextAddress | dict | list | tuple
    original_value: str
    new_value: str


@dataclass(frozen=True)
class AppliedChange:
    address: TextAddress
    original_value: str
    new_value: str
    success: bool
    reason: str | None = None


ADDRESS_FIELDS = ["slide_idx", "shape_id", "kind", "row", "col", "paragraph_idx", "run_idx"]


def _address_from_any(value: TextAddress | dict | list | tuple) -> TextAddress:
    if isinstance(value, TextAddress):
        return value
    if is_dataclass(value):
        return TextAddress(**asdict(value))
    if isinstance(value, dict):
        return TextAddress(**{field: value.get(field) for field in ADDRESS_FIELDS})
    if isinstance(value, (list, tuple)):
        return TextAddress(*value)
    raise TypeError(f"Unsupported address type: {type(value).__name__}")


def _shape_by_id(slide, shape_id: int):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    return None


def _run_for_address(prs, address: TextAddress):
    try:
        slide = prs.slides[address.slide_idx]
    except IndexError:
        return None, "slide not found"

    shape = _shape_by_id(slide, address.shape_id)
    if shape is None:
        return None, "shape not found"

    if address.kind == "text_frame":
        if not getattr(shape, "has_text_frame", False):
            return None, "shape has no text frame"
        paragraphs = shape.text_frame.paragraphs
    elif address.kind == "table_cell":
        if not getattr(shape, "has_table", False):
            return None, "shape has no table"
        if address.row is None or address.col is None:
            return None, "table cell address missing row or column"
        try:
            paragraphs = shape.table.cell(address.row, address.col).text_frame.paragraphs
        except (IndexError, ValueError):
            return None, "table cell not found"
    else:
        return None, "unsupported address kind"

    try:
        paragraph = paragraphs[address.paragraph_idx]
    except IndexError:
        return None, "paragraph not found"

    try:
        return paragraph.runs[address.run_idx], None
    except IndexError:
        return None, "run not found"


def _mapping_value(mapping: Mapping | dict, name: str) -> Any:
    if isinstance(mapping, dict):
        return mapping[name]
    return getattr(mapping, name)


def apply_token_mappings(file_bytes: bytes, mappings: list[Mapping]) -> tuple[bytes, list[AppliedChange]]:
    """Apply deterministic single-run token replacements to a PPTX.

    TODO(RR-1 Phase C): if PowerPoint splits a visible token across multiple
    runs, detect it during ingest and rewrite the paragraph using the first
    run's font properties while dropping the continuation runs. Phase A ships
    the simple single-run replacement path only.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    changes: list[AppliedChange] = []

    for mapping in mappings:
        address = _address_from_any(_mapping_value(mapping, "address"))
        original_value = str(_mapping_value(mapping, "original_value"))
        new_value = str(_mapping_value(mapping, "new_value"))

        run, reason = _run_for_address(prs, address)
        if run is None:
            changes.append(AppliedChange(address, original_value, new_value, False, reason))
            continue
        if original_value not in run.text:
            changes.append(
                AppliedChange(address, original_value, new_value, False, "token not found in run")
            )
            continue

        run.text = run.text.replace(original_value, new_value, 1)
        changes.append(AppliedChange(address, original_value, new_value, True, None))

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue(), changes
