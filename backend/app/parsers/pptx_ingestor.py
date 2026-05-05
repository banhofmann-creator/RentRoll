from __future__ import annotations

import io
import re
from dataclasses import asdict, dataclass
from typing import Literal

from pptx import Presentation

from app.core.kpi_catalog import KPI_CATALOG

TOKEN_RE = re.compile(r"\{\{([a-zA-Z0-9_]+)\}\}")


@dataclass(frozen=True)
class TextAddress:
    slide_idx: int
    shape_id: int
    kind: Literal["text_frame", "table_cell"]
    row: int | None
    col: int | None
    paragraph_idx: int
    run_idx: int


@dataclass(frozen=True)
class TextElement:
    address: TextAddress
    text: str
    font_name: str | None
    font_size_pt: float | None
    font_bold: bool | None
    font_color_rgb: str | None


@dataclass(frozen=True)
class TokenCandidate:
    address: TextAddress
    kpi_id: str
    full_text: str
    span: tuple[int, int]


def _font_color_rgb(run) -> str | None:
    try:
        rgb = run.font.color.rgb
    except Exception:
        return None
    return str(rgb) if rgb is not None else None


def _run_to_element(
    slide_idx: int,
    shape_id: int,
    kind: Literal["text_frame", "table_cell"],
    row: int | None,
    col: int | None,
    paragraph_idx: int,
    run_idx: int,
    run,
) -> TextElement:
    size = run.font.size
    return TextElement(
        address=TextAddress(
            slide_idx=slide_idx,
            shape_id=shape_id,
            kind=kind,
            row=row,
            col=col,
            paragraph_idx=paragraph_idx,
            run_idx=run_idx,
        ),
        text=run.text,
        font_name=run.font.name,
        font_size_pt=float(size.pt) if size is not None else None,
        font_bold=run.font.bold,
        font_color_rgb=_font_color_rgb(run),
    )


def ingest_pptx(file_bytes: bytes) -> list[TextElement]:
    prs = Presentation(io.BytesIO(file_bytes))
    elements: list[TextElement] = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        elements.append(
                            _run_to_element(
                                slide_idx,
                                shape.shape_id,
                                "text_frame",
                                None,
                                None,
                                paragraph_idx,
                                run_idx,
                                run,
                            )
                        )

            if getattr(shape, "has_table", False):
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        for paragraph_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                            for run_idx, run in enumerate(paragraph.runs):
                                elements.append(
                                    _run_to_element(
                                        slide_idx,
                                        shape.shape_id,
                                        "table_cell",
                                        row_idx,
                                        col_idx,
                                        paragraph_idx,
                                        run_idx,
                                        run,
                                    )
                                )

    return elements


def find_token_candidates(elements: list[TextElement]) -> tuple[list[TokenCandidate], list[str]]:
    candidates: list[TokenCandidate] = []
    unknown_tokens: list[str] = []

    for element in elements:
        for match in TOKEN_RE.finditer(element.text):
            kpi_id = match.group(1)
            if kpi_id not in KPI_CATALOG:
                if kpi_id not in unknown_tokens:
                    unknown_tokens.append(kpi_id)
                continue
            candidates.append(
                TokenCandidate(
                    address=element.address,
                    kpi_id=kpi_id,
                    full_text=match.group(0),
                    span=match.span(),
                )
            )

    return candidates, unknown_tokens


def token_candidate_to_dict(candidate: TokenCandidate) -> dict:
    data = asdict(candidate)
    data["span"] = list(candidate.span)
    return data
